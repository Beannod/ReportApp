from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

names_text = '''Sanstha Ko Logo(Tribeni Youth )
8'Th Program
Thnakyou So Much For All The Helping Hand
Rajat Kafle
Pratik Bashyal
Hotel Grand Tribeni
Manish Gurung
Binod Bhandari
Victor Shahi Thakuri
Sudip Rayamajhi
Sameer Adhikari,Sudin Sharma,Kshitiz Sahani
Bishwo Karki
Saroj Ale Magar
Madhu Sudan Rokka Chettri(Madan Dai)
Krishna Poudel
Akash Sharma
Bishal Gharti Magar
Tanka Shrestha
Shova Shrestha
Bibek Thapa
Tribeni Youth Vision
Sunil Shrestha
Ajit Gurung
Bir Bahadur Singh
Sudeep Chhetri
Kamal Thapa
Amrit Shrestha
Subham Chettri
Kabita Shrestha
Binod Dhakal
Sanjog Thapa
Raj Kumar Rai
Bhoj Bahadur Thapa
Rajiv Gurung
Tara Kunwar
Rabindra Sapkota
Shree Shree Acharya Chakrapani Swami
Diyo Maya Pun
Gs Cafe Family Restaurant And Lodge(Propriter/Bishwo Mahardhan Shah)
Triveni View Point
Ajay Baniya
Greenland Hotel
Khusbu Fancy Store(With Photo)
Shubash Chettri
Prakash Gurung
Janak Bhandari
Sapana Gurung
Kamal Thapa
Sagar Thapa
Manni Thapa
Shrijana Sampang Rai
Chadra Bahadur Bista
Lalji Kurmi
Keshav Gautam
Rishi Shrestha
Sandeep Nepal
Mukti Thapa
Nirmal Gurung
Anish Sharma
Sagar Rana
Barsat Thapa Magar
Tara Gurung
Bijay Gaud
Dilip Thapa
Dilip Sharma
Shree Ram Gaudel
Damodar Poudel
Ramesh Chaurasiya
Anil Pun
Aklesh Chaurasiya
Ishwor Thakuri
Tara Bahdaur Karki
Sanjay Kc
Moti Lal Upadhaya
Santosh Gurung
Prabin Rana Magar
Nirmal Sunar
Ganga Bahadur Pun
Yuvraj Thakuri'''

def make_slide(out_path='Tribeni_Youth_Sponsors.pptx'):
    prs = Presentation()
    # blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    p.text = 'Tribeni Youth — Program Supporters & Helpers'
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # Prepare names list
    lines = [ln.strip() for ln in names_text.splitlines() if ln.strip()]
    # split into two roughly equal columns
    mid = (len(lines) + 1) // 2
    col1 = lines[:mid]
    col2 = lines[mid:]

    left = Inches(0.5)
    top = Inches(1.4)
    width = Inches(4.25)
    height = Inches(6.0)

    tx1 = slide.shapes.add_textbox(left, top, width, height)
    tf1 = tx1.text_frame
    tf1.margin_left = Pt(6)
    tf1.margin_right = Pt(6)
    for i, name in enumerate(col1):
        p = tf1.add_paragraph() if i > 0 else tf1.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)

    tx2 = slide.shapes.add_textbox(left + width + Inches(0.2), top, width, height)
    tf2 = tx2.text_frame
    tf2.margin_left = Pt(6)
    tf2.margin_right = Pt(6)
    for i, name in enumerate(col2):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)

    # small footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(8.0), Inches(9), Inches(0.4))
    ftf = footer.text_frame
    fp = ftf.paragraphs[0]
    fp.text = 'Thank you to everyone who helped organize this cultural program.'
    fp.font.size = Pt(10)
    fp.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    prs.save(out_path)
    print(f'Wrote {out_path}')

if __name__ == '__main__':
    make_slide()

