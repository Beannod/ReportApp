#!/usr/bin/env python3
"""Convert a directory of images into a single PowerPoint (.pptx).

Usage:
    python scripts\images_to_pptx.py --input-dir path\to\images --output-file output.pptx

The script will add one image per slide and scale each image to fit the slide while
preserving aspect ratio and centering it.
"""
import argparse
import os
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt


IMAGE_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'}


def images_in_dir(directory: Path):
    files = [p for p in sorted(directory.iterdir()) if p.suffix.lower() in IMAGE_EXTS and p.is_file()]
    return files


def add_image_slide(prs: Presentation, image_path: Path):
    # blank layout
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # slide dimensions in inches
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # open image to get size
    with Image.open(image_path) as img:
        width_px, height_px = img.size

    # pptx uses EMU units; we can compute target size preserving aspect ratio
    # We'll fit the image within the slide margins (small margin)
    margin = Inches(0.5)
    max_w = slide_w - margin * 2
    max_h = slide_h - margin * 2

    # convert px to EMU requires DPI; python-pptx add_picture accepts inches or emu via Inches()
    # We'll compute target inches from pixel size using a fallback DPI of 96
    dpi = 96.0
    img_w_in = width_px / dpi
    img_h_in = height_px / dpi

    # Convert slide EMU to inches for comparison
    from pptx.util import Emu
    slide_w_in = slide_w / Emu(1)
    slide_h_in = slide_h / Emu(1)

    max_w_in = max_w / Emu(1)
    max_h_in = max_h / Emu(1)

    ratio = min(max_w_in / img_w_in, max_h_in / img_h_in, 1.0)
    target_w_in = img_w_in * ratio
    target_h_in = img_h_in * ratio

    # center position
    left = (slide_w_in - target_w_in) / 2
    top = (slide_h_in - target_h_in) / 2

    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(target_w_in), height=Inches(target_h_in))


def make_pptx(input_dir: Path, output_file: Path):
    imgs = images_in_dir(input_dir)
    if not imgs:
        raise SystemExit(f'No images found in {input_dir!s}')

    prs = Presentation()

    for img in imgs:
        add_image_slide(prs, img)

    prs.save(str(output_file))
    print(f'Wrote {output_file!s} with {len(imgs)} slides')


def main():
    p = argparse.ArgumentParser(description='Create a PPTX from images in a folder')
    p.add_argument('--input-dir', '-i', required=True, help='Directory containing images')
    p.add_argument('--output-file', '-o', default='images.pptx', help='Output PPTX filename')
    args = p.parse_args()

    inp = Path(args.input_dir)
    if not inp.exists() or not inp.is_dir():
        raise SystemExit(f'Input directory not found: {inp!s}')

    out = Path(args.output_file)
    make_pptx(inp, out)


if __name__ == '__main__':
    main()
